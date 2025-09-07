"""
Utility functions for Format Service
"""

import os
import re
import magic
import logging
import asyncio
from typing import List, Optional, Any
from pathlib import Path
import tempfile
import shutil
from datetime import datetime, timedelta

# Document processing imports
import docx
from PyPDF2 import PdfReader
import pdfplumber
from openpyxl import load_workbook
from pptx import Presentation
import pytesseract
from PIL import Image
import pdf2image

logger = logging.getLogger(__name__)

def detect_file_type(file_path: str) -> str:
    """Detect file type using python-magic"""
    try:
        mime_type = magic.from_file(file_path, mime=True)
        
        # Map MIME types to file extensions
        mime_to_ext = {
            'application/pdf': 'pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
            'application/msword': 'doc',
            'text/plain': 'txt',
            'application/rtf': 'rtf',
            'application/vnd.oasis.opendocument.text': 'odt',
            'text/html': 'html',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
            'application/vnd.ms-excel': 'xls',
            'application/vnd.oasis.opendocument.spreadsheet': 'ods',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
            'application/vnd.ms-powerpoint': 'ppt',
            'application/vnd.oasis.opendocument.presentation': 'odp'
        }
        
        detected_type = mime_to_ext.get(mime_type)
        
        if not detected_type:
            # Fallback to file extension
            ext = Path(file_path).suffix[1:].lower()
            if ext in ['pdf', 'docx', 'doc', 'txt', 'rtf', 'odt', 'html', 'xlsx', 'xls', 'ods', 'pptx', 'ppt', 'odp']:
                detected_type = ext
            else:
                detected_type = 'unknown'
        
        return detected_type
        
    except Exception as e:
        logger.error(f"File type detection failed: {e}")
        # Fallback to extension
        return Path(file_path).suffix[1:].lower() or 'unknown'

def validate_file_size(file, max_size_mb: int = 100) -> bool:
    """Validate file size"""
    try:
        # Get file size
        file.file.seek(0, 2)  # Seek to end
        size = file.file.tell()
        file.file.seek(0)  # Reset position
        
        max_size_bytes = max_size_mb * 1024 * 1024
        return size <= max_size_bytes
        
    except Exception as e:
        logger.error(f"File size validation failed: {e}")
        return False

def extract_text_content(file_path: str, file_type: str, max_chars: Optional[int] = None) -> str:
    """Extract text content from various file types"""
    try:
        text_content = ""
        
        if file_type == 'pdf':
            text_content = extract_text_from_pdf(file_path, max_chars)
        elif file_type == 'docx':
            text_content = extract_text_from_docx(file_path, max_chars)
        elif file_type == 'doc':
            text_content = extract_text_from_doc(file_path, max_chars)
        elif file_type in ['xlsx', 'xls']:
            text_content = extract_text_from_excel(file_path, max_chars)
        elif file_type in ['pptx', 'ppt']:
            text_content = extract_text_from_pptx(file_path, max_chars)
        elif file_type == 'txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read()
        elif file_type == 'rtf':
            text_content = extract_text_from_rtf(file_path, max_chars)
        elif file_type == 'html':
            text_content = extract_text_from_html(file_path, max_chars)
        
        # Truncate if max_chars specified
        if max_chars and len(text_content) > max_chars:
            text_content = text_content[:max_chars] + "..."
        
        return text_content
        
    except Exception as e:
        logger.error(f"Text extraction failed for {file_type}: {e}")
        return ""

def extract_text_from_pdf(file_path: str, max_chars: Optional[int] = None) -> str:
    """Extract text from PDF"""
    try:
        text_content = ""
        
        # Try pdfplumber first (better for complex layouts)
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content += page_text + "\n\n"
                
                # Check if we've reached max_chars
                if max_chars and len(text_content) > max_chars:
                    break
        
        # If no text found, try OCR
        if not text_content.strip():
            text_content = extract_text_from_pdf_ocr(file_path, max_chars)
        
        return text_content
        
    except Exception as e:
        logger.error(f"PDF text extraction failed: {e}")
        # Fallback to PyPDF2
        try:
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                text_parts = []
                for page in reader.pages:
                    text_parts.append(page.extract_text())
                    if max_chars and len('\n'.join(text_parts)) > max_chars:
                        break
                return '\n'.join(text_parts)
        except:
            return ""

def extract_text_from_pdf_ocr(file_path: str, max_chars: Optional[int] = None) -> str:
    """Extract text from PDF using OCR"""
    try:
        # Convert PDF pages to images
        images = pdf2image.convert_from_path(file_path, dpi=200)
        
        text_content = ""
        for i, image in enumerate(images[:5]):  # Process max 5 pages for OCR
            # Perform OCR
            page_text = pytesseract.image_to_string(image)
            text_content += page_text + "\n\n"
            
            if max_chars and len(text_content) > max_chars:
                break
        
        return text_content
        
    except Exception as e:
        logger.error(f"PDF OCR extraction failed: {e}")
        return ""

def extract_text_from_docx(file_path: str, max_chars: Optional[int] = None) -> str:
    """Extract text from DOCX"""
    try:
        doc = docx.Document(file_path)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
                
                if max_chars and len('\n'.join(text_parts)) > max_chars:
                    break
        
        return '\n'.join(text_parts)
        
    except Exception as e:
        logger.error(f"DOCX text extraction failed: {e}")
        return ""

def extract_text_from_doc(file_path: str, max_chars: Optional[int] = None) -> str:
    """Extract text from legacy DOC files"""
    try:
        # Try using antiword if available
        import subprocess
        result = subprocess.run(['antiword', file_path], capture_output=True, text=True)
        if result.returncode == 0:
            text = result.stdout
            if max_chars and len(text) > max_chars:
                text = text[:max_chars]
            return text
        else:
            logger.warning(f"Antiword failed for {file_path}")
            return ""
            
    except FileNotFoundError:
        logger.warning("Antiword not available for DOC extraction")
        return ""
    except Exception as e:
        logger.error(f"DOC text extraction failed: {e}")
        return ""

def extract_text_from_excel(file_path: str, max_chars: Optional[int] = None) -> str:
    """Extract text from Excel files"""
    try:
        wb = load_workbook(file_path, read_only=True)
        text_parts = []
        
        for sheet in wb.worksheets:
            text_parts.append(f"Sheet: {sheet.title}")
            
            for row in sheet.iter_rows():
                row_data = []
                for cell in row:
                    if cell.value:
                        row_data.append(str(cell.value))
                
                if row_data:
                    text_parts.append('\t'.join(row_data))
                
                # Check max_chars
                current_text = '\n'.join(text_parts)
                if max_chars and len(current_text) > max_chars:
                    return current_text[:max_chars]
            
            text_parts.append("")  # Empty line between sheets
        
        return '\n'.join(text_parts)
        
    except Exception as e:
        logger.error(f"Excel text extraction failed: {e}")
        return ""

def extract_text_from_pptx(file_path: str, max_chars: Optional[int] = None) -> str:
    """Extract text from PowerPoint files"""
    try:
        prs = Presentation(file_path)
        text_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"Slide {i}:")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
            
            text_parts.append("")  # Empty line between slides
            
            # Check max_chars
            current_text = '\n'.join(text_parts)
            if max_chars and len(current_text) > max_chars:
                return current_text[:max_chars]
        
        return '\n'.join(text_parts)
        
    except Exception as e:
        logger.error(f"PPTX text extraction failed: {e}")
        return ""

def extract_text_from_rtf(file_path: str, max_chars: Optional[int] = None) -> str:
    """Extract text from RTF files"""
    try:
        # Simple RTF text extraction (basic implementation)
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        # Remove RTF control codes (very basic)
        text = re.sub(r'\\[a-z]+\d*', '', content)
        text = re.sub(r'[{}]', '', text)
        text = re.sub(r'\\', '', text)
        text = ' '.join(text.split())
        
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
        
        return text
        
    except Exception as e:
        logger.error(f"RTF text extraction failed: {e}")
        return ""

def extract_text_from_html(file_path: str, max_chars: Optional[int] = None) -> str:
    """Extract text from HTML files"""
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        text = soup.get_text()
        
        # Clean up whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        
        if max_chars and len(text) > max_chars:
            text = text[:max_chars]
        
        return text
        
    except Exception as e:
        logger.error(f"HTML text extraction failed: {e}")
        return ""

def segment_text(text: str, options: Optional[Any] = None) -> List[str]:
    """Segment text for translation"""
    try:
        if not text.strip():
            return []
        
        # Default segmentation by sentences
        sentences = segment_by_sentences(text)
        
        if options:
            if hasattr(options, 'rule'):
                if options.rule == 'paragraph':
                    sentences = segment_by_paragraphs(text)
                elif options.rule == 'line':
                    sentences = text.split('\n')
                elif options.rule == 'custom' and hasattr(options, 'custom_patterns'):
                    sentences = segment_by_custom_patterns(text, options.custom_patterns)
            
            # Apply length constraints
            if hasattr(options, 'min_segment_length') or hasattr(options, 'max_segment_length'):
                min_len = getattr(options, 'min_segment_length', 10)
                max_len = getattr(options, 'max_segment_length', 500)
                sentences = apply_length_constraints(sentences, min_len, max_len)
            
            # Merge short segments if requested
            if hasattr(options, 'merge_short_segments') and options.merge_short_segments:
                sentences = merge_short_segments(sentences, 50)  # Default threshold
        
        # Clean and filter segments
        cleaned_segments = []
        for segment in sentences:
            cleaned = segment.strip()
            if cleaned:
                cleaned_segments.append(cleaned)
        
        return cleaned_segments
        
    except Exception as e:
        logger.error(f"Text segmentation failed: {e}")
        return [text]  # Return original as single segment

def segment_by_sentences(text: str) -> List[str]:
    """Segment text by sentences"""
    # Simple sentence boundary detection
    sentence_endings = r'[.!?]+\s+'
    sentences = re.split(sentence_endings, text)
    
    # Clean up sentences
    cleaned = []
    for sentence in sentences:
        sentence = sentence.strip()
        if sentence:
            cleaned.append(sentence)
    
    return cleaned

def segment_by_paragraphs(text: str) -> List[str]:
    """Segment text by paragraphs"""
    paragraphs = text.split('\n\n')
    
    cleaned = []
    for para in paragraphs:
        para = para.strip()
        if para:
            cleaned.append(para)
    
    return cleaned

def segment_by_custom_patterns(text: str, patterns: List[str]) -> List[str]:
    """Segment text using custom regex patterns"""
    try:
        if not patterns:
            return segment_by_sentences(text)
        
        # Use first pattern for segmentation
        pattern = patterns[0]
        segments = re.split(pattern, text)
        
        cleaned = []
        for segment in segments:
            segment = segment.strip()
            if segment:
                cleaned.append(segment)
        
        return cleaned
        
    except Exception as e:
        logger.error(f"Custom pattern segmentation failed: {e}")
        return segment_by_sentences(text)

def apply_length_constraints(segments: List[str], min_len: int, max_len: int) -> List[str]:
    """Apply minimum and maximum length constraints to segments"""
    constrained = []
    
    for segment in segments:
        if len(segment) < min_len:
            # Skip very short segments or merge with next
            if constrained:
                constrained[-1] += " " + segment
            else:
                constrained.append(segment)
        elif len(segment) > max_len:
            # Split long segments
            parts = split_long_segment(segment, max_len)
            constrained.extend(parts)
        else:
            constrained.append(segment)
    
    return constrained

def split_long_segment(segment: str, max_len: int) -> List[str]:
    """Split long segment into smaller parts"""
    parts = []
    words = segment.split()
    current_part = []
    current_length = 0
    
    for word in words:
        word_length = len(word) + 1  # +1 for space
        
        if current_length + word_length <= max_len:
            current_part.append(word)
            current_length += word_length
        else:
            if current_part:
                parts.append(' '.join(current_part))
                current_part = [word]
                current_length = len(word)
            else:
                # Single word is too long, add it anyway
                parts.append(word)
    
    if current_part:
        parts.append(' '.join(current_part))
    
    return parts

def merge_short_segments(segments: List[str], threshold: int = 50) -> List[str]:
    """Merge segments that are shorter than threshold"""
    if not segments:
        return segments
    
    merged = []
    current_segment = segments[0]
    
    for i in range(1, len(segments)):
        next_segment = segments[i]
        
        if len(current_segment) < threshold:
            # Merge with next segment
            current_segment += " " + next_segment
        else:
            merged.append(current_segment)
            current_segment = next_segment
    
    # Add the last segment
    merged.append(current_segment)
    
    return merged

def normalize_text(text: str) -> str:
    """Normalize text for processing"""
    if not text:
        return text
    
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text.strip())
    
    # Fix common OCR issues
    text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)  # Split camelCase
    text = re.sub(r'(\d)([A-Z])', r'\1 \2', text)     # Split numberLetter
    text = re.sub(r'([a-z])(\d)', r'\1 \2', text)     # Split letterNumber
    
    # Normalize punctuation spacing
    text = re.sub(r'\s*([,.;:!?])\s*', r'\1 ', text)
    text = re.sub(r'\s+', ' ', text.strip())
    
    return text

async def cleanup_temp_files(directory: str, delay_minutes: int = 60):
    """Cleanup temporary files after delay"""
    try:
        # Wait for the specified delay
        await asyncio.sleep(delay_minutes * 60)
        
        # Remove directory if it exists
        if os.path.exists(directory):
            shutil.rmtree(directory)
            logger.info(f"Cleaned up temporary directory: {directory}")
            
    except Exception as e:
        logger.error(f"Cleanup failed for {directory}: {e}")

def cleanup_old_files(base_dir: str, max_age_hours: int = 24):
    """Cleanup files older than specified age"""
    try:
        cutoff_time = datetime.now() - timedelta(hours=max_age_hours)
        
        for root, dirs, files in os.walk(base_dir):
            for file in files:
                file_path = os.path.join(root, file)
                
                try:
                    file_time = datetime.fromtimestamp(os.path.getmtime(file_path))
                    if file_time < cutoff_time:
                        os.unlink(file_path)
                        logger.debug(f"Cleaned up old file: {file_path}")
                except Exception as e:
                    logger.error(f"Failed to cleanup {file_path}: {e}")
            
            # Remove empty directories
            for dir_name in dirs:
                dir_path = os.path.join(root, dir_name)
                try:
                    if not os.listdir(dir_path):
                        os.rmdir(dir_path)
                        logger.debug(f"Removed empty directory: {dir_path}")
                except Exception as e:
                    logger.error(f"Failed to remove directory {dir_path}: {e}")
                    
    except Exception as e:
        logger.error(f"Old file cleanup failed: {e}")