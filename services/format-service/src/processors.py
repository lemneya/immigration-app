"""
Document processors for format conversion and XLIFF generation
"""

import os
import subprocess
import asyncio
import logging
import xml.etree.ElementTree as ET
from typing import Dict, List, Optional, Any
from datetime import datetime
import tempfile
import shutil
import json
from pathlib import Path

# Document processing imports
import docx
from PyPDF2 import PdfReader
import pdfplumber
from openpyxl import load_workbook
from pptx import Presentation

# Text processing
import re
from langdetect import detect
from .utils import segment_text, extract_text_content, normalize_text

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handle document format conversion"""
    
    def __init__(self):
        self.supported_input = ['pdf', 'docx', 'doc', 'txt', 'rtf', 'odt', 'xlsx', 'xls', 'pptx', 'ppt']
        self.supported_output = ['txt', 'docx', 'pdf', 'html']
    
    async def analyze_document(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """Analyze document and extract metadata"""
        try:
            info = {
                'file_type': file_type,
                'word_count': 0,
                'page_count': 1,
                'language': 'unknown',
                'metadata': {},
                'extractable_formats': []
            }
            
            if file_type == 'pdf':
                info.update(await self._analyze_pdf(file_path))
            elif file_type in ['docx', 'doc']:
                info.update(await self._analyze_docx(file_path))
            elif file_type in ['xlsx', 'xls']:
                info.update(await self._analyze_excel(file_path))
            elif file_type in ['pptx', 'ppt']:
                info.update(await self._analyze_pptx(file_path))
            elif file_type == 'txt':
                info.update(await self._analyze_text(file_path))
            
            # Determine extractable formats
            info['extractable_formats'] = self._get_extractable_formats(file_type)
            
            return info
            
        except Exception as e:
            logger.error(f"Document analysis failed: {e}")
            return {'error': str(e)}
    
    async def _analyze_pdf(self, file_path: str) -> Dict[str, Any]:
        """Analyze PDF document"""
        try:
            with pdfplumber.open(file_path) as pdf:
                page_count = len(pdf.pages)
                
                # Extract text from first few pages for analysis
                text_sample = ""
                for i, page in enumerate(pdf.pages[:3]):  # First 3 pages
                    text_sample += page.extract_text() or ""
                
                word_count = len(text_sample.split()) if text_sample else 0
                
                # Detect language
                language = 'unknown'
                if text_sample and len(text_sample.strip()) > 50:
                    try:
                        language = detect(text_sample[:1000])
                    except:
                        pass
                
                return {
                    'page_count': page_count,
                    'word_count': word_count,
                    'language': language,
                    'metadata': {
                        'has_text': bool(text_sample.strip()),
                        'extractable_text': len(text_sample) > 100
                    }
                }
                
        except Exception as e:
            logger.error(f"PDF analysis failed: {e}")
            return {'error': str(e)}
    
    async def _analyze_docx(self, file_path: str) -> Dict[str, Any]:
        """Analyze Word document"""
        try:
            if file_path.endswith('.docx'):
                doc = docx.Document(file_path)
                
                # Extract text
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                text_content = '\n'.join(full_text)
                
                word_count = len(text_content.split()) if text_content else 0
                
                # Detect language
                language = 'unknown'
                if text_content and len(text_content.strip()) > 50:
                    try:
                        language = detect(text_content[:1000])
                    except:
                        pass
                
                return {
                    'page_count': 1,  # Approximate
                    'word_count': word_count,
                    'language': language,
                    'metadata': {
                        'paragraphs': len(doc.paragraphs),
                        'has_images': len([r for r in doc.inline_shapes]) > 0
                    }
                }
            else:
                # For .doc files, use text extraction
                return await self._analyze_legacy_doc(file_path)
                
        except Exception as e:
            logger.error(f"DOCX analysis failed: {e}")
            return {'error': str(e)}
    
    async def _analyze_legacy_doc(self, file_path: str) -> Dict[str, Any]:
        """Analyze legacy .doc file"""
        try:
            # Use external tool like antiword or catdoc if available
            # For now, return basic info
            return {
                'page_count': 1,
                'word_count': 0,
                'language': 'unknown',
                'metadata': {'legacy_format': True}
            }
        except Exception as e:
            return {'error': str(e)}
    
    async def _analyze_excel(self, file_path: str) -> Dict[str, Any]:
        """Analyze Excel document"""
        try:
            wb = load_workbook(file_path, read_only=True)
            
            total_cells = 0
            text_content = []
            
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            total_cells += 1
                            if isinstance(cell.value, str):
                                text_content.append(cell.value)
            
            full_text = ' '.join(text_content)
            word_count = len(full_text.split()) if full_text else 0
            
            # Detect language
            language = 'unknown'
            if full_text and len(full_text.strip()) > 50:
                try:
                    language = detect(full_text[:1000])
                except:
                    pass
            
            return {
                'page_count': len(wb.worksheets),
                'word_count': word_count,
                'language': language,
                'metadata': {
                    'sheets': len(wb.worksheets),
                    'cells_with_data': total_cells
                }
            }
            
        except Exception as e:
            logger.error(f"Excel analysis failed: {e}")
            return {'error': str(e)}
    
    async def _analyze_pptx(self, file_path: str) -> Dict[str, Any]:
        """Analyze PowerPoint document"""
        try:
            prs = Presentation(file_path)
            
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            
            full_text = ' '.join(text_content)
            word_count = len(full_text.split()) if full_text else 0
            
            # Detect language
            language = 'unknown'
            if full_text and len(full_text.strip()) > 50:
                try:
                    language = detect(full_text[:1000])
                except:
                    pass
            
            return {
                'page_count': len(prs.slides),
                'word_count': word_count,
                'language': language,
                'metadata': {
                    'slides': len(prs.slides)
                }
            }
            
        except Exception as e:
            logger.error(f"PPTX analysis failed: {e}")
            return {'error': str(e)}
    
    async def _analyze_text(self, file_path: str) -> Dict[str, Any]:
        """Analyze plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            word_count = len(content.split()) if content else 0
            line_count = len(content.splitlines())
            
            # Detect language
            language = 'unknown'
            if content and len(content.strip()) > 50:
                try:
                    language = detect(content[:1000])
                except:
                    pass
            
            return {
                'page_count': max(1, line_count // 50),  # Approximate
                'word_count': word_count,
                'language': language,
                'metadata': {
                    'line_count': line_count,
                    'char_count': len(content)
                }
            }
            
        except Exception as e:
            logger.error(f"Text analysis failed: {e}")
            return {'error': str(e)}
    
    def _get_extractable_formats(self, input_format: str) -> List[str]:
        """Get list of formats this input can be converted to"""
        base_formats = ['txt']
        
        if input_format in ['pdf', 'docx']:
            base_formats.extend(['html', 'docx'])
        elif input_format in ['xlsx', 'xls']:
            base_formats.extend(['csv', 'html'])
        elif input_format in ['pptx', 'ppt']:
            base_formats.extend(['html'])
        
        return base_formats
    
    async def convert_document(self, input_path: str, target_format: str, 
                             output_dir: str, options: Optional[Any] = None) -> Dict[str, Any]:
        """Convert document to target format"""
        try:
            input_format = Path(input_path).suffix[1:].lower()
            base_name = Path(input_path).stem
            output_filename = f"{base_name}.{target_format}"
            output_path = os.path.join(output_dir, output_filename)
            
            if target_format == 'txt':
                # Extract text content
                text_content = extract_text_content(input_path, input_format)
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text_content)
            
            elif target_format == 'html' and input_format == 'docx':
                # Convert DOCX to HTML
                await self._convert_docx_to_html(input_path, output_path)
            
            elif target_format == 'docx' and input_format == 'txt':
                # Convert text to DOCX
                await self._convert_txt_to_docx(input_path, output_path)
            
            else:
                raise ValueError(f"Conversion from {input_format} to {target_format} not supported")
            
            # Calculate word count
            if target_format == 'txt':
                with open(output_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                word_count = len(content.split())
            else:
                word_count = None
            
            return {
                'filename': output_filename,
                'path': output_path,
                'word_count': word_count,
                'metadata': {
                    'input_format': input_format,
                    'output_format': target_format
                }
            }
            
        except Exception as e:
            logger.error(f"Document conversion failed: {e}")
            raise
    
    async def _convert_docx_to_html(self, input_path: str, output_path: str):
        """Convert DOCX to HTML"""
        try:
            doc = docx.Document(input_path)
            
            html_content = ['<html><body>']
            for para in doc.paragraphs:
                if para.text.strip():
                    html_content.append(f'<p>{para.text}</p>')
            html_content.append('</body></html>')
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(html_content))
                
        except Exception as e:
            logger.error(f"DOCX to HTML conversion failed: {e}")
            raise
    
    async def _convert_txt_to_docx(self, input_path: str, output_path: str):
        """Convert text to DOCX"""
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            doc = docx.Document()
            
            # Split content into paragraphs
            paragraphs = content.split('\n\n')
            for para_text in paragraphs:
                if para_text.strip():
                    doc.add_paragraph(para_text.strip())
            
            doc.save(output_path)
            
        except Exception as e:
            logger.error(f"Text to DOCX conversion failed: {e}")
            raise

class XLIFFProcessor:
    """Handle XLIFF creation and processing"""
    
    def __init__(self):
        self.xliff_version = "2.1"
    
    async def create_xliff(self, input_path: str, source_lang: str, target_lang: str,
                          output_dir: str, options: Optional[Any] = None) -> Dict[str, Any]:
        """Create XLIFF from document"""
        try:
            # Extract text content
            file_type = Path(input_path).suffix[1:].lower()
            text_content = extract_text_content(input_path, file_type)
            
            if not text_content.strip():
                raise ValueError("No extractable text found in document")
            
            # Segment text
            segments = segment_text(text_content, options)
            
            # Create XLIFF
            base_name = Path(input_path).stem
            xliff_filename = f"{base_name}_{source_lang}_{target_lang}.xliff"
            xliff_path = os.path.join(output_dir, xliff_filename)
            
            await self._create_xliff_file(
                segments=segments,
                source_lang=source_lang,
                target_lang=target_lang,
                output_path=xliff_path,
                original_filename=Path(input_path).name
            )
            
            # Calculate statistics
            word_count = sum(len(seg.split()) for seg in segments)
            
            # Create preview segments
            preview_segments = []
            for i, seg in enumerate(segments[:5]):  # First 5 segments
                preview_segments.append({
                    'id': f"seg_{i+1}",
                    'source_text': seg,
                    'target_text': "",
                    'status': "new",
                    'approved': False,
                    'locked': False,
                    'notes': [],
                    'metadata': {}
                })
            
            return {
                'xliff_filename': xliff_filename,
                'segment_count': len(segments),
                'word_count': word_count,
                'preview_segments': preview_segments,
                'metadata': {
                    'xliff_version': self.xliff_version,
                    'source_language': source_lang,
                    'target_language': target_lang,
                    'original_filename': Path(input_path).name
                }
            }
            
        except Exception as e:
            logger.error(f"XLIFF creation failed: {e}")
            raise
    
    async def _create_xliff_file(self, segments: List[str], source_lang: str, 
                                target_lang: str, output_path: str, original_filename: str):
        """Create XLIFF 2.1 file"""
        try:
            # Create XLIFF structure
            xliff_content = [
                '<?xml version="1.0" encoding="UTF-8"?>',
                f'<xliff xmlns="urn:oasis:names:tc:xliff:document:2.1" version="2.1" srcLang="{source_lang}" trgLang="{target_lang}">',
                f'  <file id="f1" original="{original_filename}">',
                '    <unit>'
            ]
            
            # Add segments
            for i, segment in enumerate(segments, 1):
                segment_id = f"seg_{i}"
                escaped_segment = self._escape_xml(segment)
                
                xliff_content.extend([
                    f'      <segment id="{segment_id}">',
                    f'        <source>{escaped_segment}</source>',
                    '        <target></target>',
                    '      </segment>'
                ])
            
            xliff_content.extend([
                '    </unit>',
                '  </file>',
                '</xliff>'
            ])
            
            # Write XLIFF file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(xliff_content))
                
        except Exception as e:
            logger.error(f"XLIFF file creation failed: {e}")
            raise
    
    def _escape_xml(self, text: str) -> str:
        """Escape XML special characters"""
        return (text.replace('&', '&amp;')
                   .replace('<', '&lt;')
                   .replace('>', '&gt;')
                   .replace('"', '&quot;')
                   .replace("'", '&#39;'))
    
    async def merge_xliff(self, xliff_path: str, output_dir: str) -> Dict[str, Any]:
        """Merge translated XLIFF back to original format"""
        try:
            # Parse XLIFF
            tree = ET.parse(xliff_path)
            root = tree.getroot()
            
            # Extract namespace
            namespace = {'xliff': 'urn:oasis:names:tc:xliff:document:2.1'}
            
            # Extract translated segments
            translated_text = []
            stats = {'total': 0, 'translated': 0, 'approved': 0}
            
            for segment in root.findall('.//xliff:segment', namespace):
                stats['total'] += 1
                
                target = segment.find('xliff:target', namespace)
                if target is not None and target.text:
                    translated_text.append(target.text)
                    stats['translated'] += 1
                else:
                    # Use source text if no translation
                    source = segment.find('xliff:source', namespace)
                    if source is not None and source.text:
                        translated_text.append(source.text)
            
            # Create output file
            base_name = Path(xliff_path).stem.replace('_en_es', '').replace('_fr_en', '').replace('_ar_en', '')
            output_filename = f"{base_name}_translated.txt"
            output_path = os.path.join(output_dir, output_filename)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(translated_text))
            
            return {
                'filename': output_filename,
                'path': output_path,
                'stats': stats
            }
            
        except Exception as e:
            logger.error(f"XLIFF merge failed: {e}")
            raise

class OkapiProcessor:
    """Handle Okapi Framework integration"""
    
    def __init__(self):
        self.okapi_root = os.getenv('OKAPI_ROOT', '/opt/okapi-apps_cocoon-m40-x64_1.44.0')
        self.tikal_path = os.path.join(self.okapi_root, 'tikal.sh')
        self.available = False
    
    async def initialize(self):
        """Initialize Okapi Framework"""
        try:
            if os.path.exists(self.tikal_path):
                # Test Okapi availability
                result = await self._run_tikal_command(['--help'])
                self.available = result['returncode'] == 0
                logger.info(f"Okapi Framework available: {self.available}")
            else:
                logger.warning(f"Okapi not found at {self.tikal_path}")
                
        except Exception as e:
            logger.error(f"Okapi initialization failed: {e}")
            self.available = False
    
    async def is_available(self) -> bool:
        """Check if Okapi is available"""
        return self.available
    
    async def is_supported_format(self, file_type: str) -> bool:
        """Check if format is supported by Okapi"""
        okapi_formats = ['docx', 'xlsx', 'pptx', 'odt', 'ods', 'odp', 'html', 'xml']
        return file_type.lower() in okapi_formats and self.available
    
    async def get_supported_formats(self) -> List[str]:
        """Get formats supported by Okapi"""
        if self.available:
            return ['docx', 'xlsx', 'pptx', 'odt', 'ods', 'odp', 'html', 'xml', 'rtf']
        return []
    
    async def create_xliff(self, input_path: str, source_lang: str, target_lang: str,
                          output_dir: str, options: Optional[Any] = None) -> Dict[str, Any]:
        """Create XLIFF using Okapi"""
        try:
            if not self.available:
                raise RuntimeError("Okapi Framework not available")
            
            # Prepare output filename
            base_name = Path(input_path).stem
            xliff_filename = f"{base_name}_{source_lang}_{target_lang}.xliff"
            xliff_path = os.path.join(output_dir, xliff_filename)
            
            # Run Okapi Tikal to extract XLIFF
            cmd = [
                '-x', input_path,          # Extract from input
                '-sl', source_lang,        # Source language
                '-tl', target_lang,        # Target language
                '-o', xliff_path          # Output XLIFF
            ]
            
            result = await self._run_tikal_command(cmd)
            
            if result['returncode'] != 0:
                raise RuntimeError(f"Okapi extraction failed: {result['stderr']}")
            
            # Parse XLIFF to get statistics
            stats = await self._analyze_xliff(xliff_path)
            
            return {
                'xliff_filename': xliff_filename,
                'segment_count': stats['segment_count'],
                'word_count': stats['word_count'],
                'preview_segments': stats['preview_segments'][:5],
                'metadata': {
                    'processor': 'okapi',
                    'xliff_version': '2.1'
                }
            }
            
        except Exception as e:
            logger.error(f"Okapi XLIFF creation failed: {e}")
            raise
    
    async def _run_tikal_command(self, args: List[str]) -> Dict[str, Any]:
        """Run Okapi Tikal command"""
        try:
            cmd = [self.tikal_path] + args
            
            process = await asyncio.create_subprocess_exec(
                *cmd,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                env={**os.environ, 'JAVA_HOME': os.getenv('JAVA_HOME', '/usr/lib/jvm/java-11-openjdk-amd64')}
            )
            
            stdout, stderr = await process.communicate()
            
            return {
                'returncode': process.returncode,
                'stdout': stdout.decode('utf-8') if stdout else '',
                'stderr': stderr.decode('utf-8') if stderr else ''
            }
            
        except Exception as e:
            logger.error(f"Tikal command failed: {e}")
            return {
                'returncode': -1,
                'stdout': '',
                'stderr': str(e)
            }
    
    async def _analyze_xliff(self, xliff_path: str) -> Dict[str, Any]:
        """Analyze XLIFF file and extract statistics"""
        try:
            tree = ET.parse(xliff_path)
            root = tree.getroot()
            
            # Find namespace
            namespace = {}
            for prefix, uri in root.nsmap.items() if hasattr(root, 'nsmap') else []:
                if uri and 'xliff' in uri:
                    namespace['xliff'] = uri
                    break
            
            if not namespace:
                # Fallback namespace detection
                namespace = {'xliff': 'urn:oasis:names:tc:xliff:document:2.1'}
            
            segments = root.findall('.//xliff:segment', namespace)
            segment_count = len(segments)
            
            word_count = 0
            preview_segments = []
            
            for i, segment in enumerate(segments[:10]):  # First 10 for preview
                source_elem = segment.find('xliff:source', namespace)
                target_elem = segment.find('xliff:target', namespace)
                
                source_text = source_elem.text if source_elem is not None else ""
                target_text = target_elem.text if target_elem is not None else ""
                
                if source_text:
                    word_count += len(source_text.split())
                
                if i < 5:  # Preview first 5
                    preview_segments.append({
                        'id': segment.get('id', f'seg_{i+1}'),
                        'source_text': source_text,
                        'target_text': target_text,
                        'status': 'new',
                        'approved': False,
                        'locked': False,
                        'notes': [],
                        'metadata': {}
                    })
            
            return {
                'segment_count': segment_count,
                'word_count': word_count,
                'preview_segments': preview_segments
            }
            
        except Exception as e:
            logger.error(f"XLIFF analysis failed: {e}")
            return {
                'segment_count': 0,
                'word_count': 0,
                'preview_segments': []
            }